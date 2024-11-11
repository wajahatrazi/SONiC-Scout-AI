
import os
import json
import openai
from PyPDF2 import PdfReader
from docx import Document
from pptx import Presentation
from dotenv import find_dotenv, load_dotenv
from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain_openai import OpenAIEmbeddings
from langchain_community.vectorstores import FAISS

# Load environment variables
dotenv_path = find_dotenv()
load_dotenv(dotenv_path)
directory = os.getenv("DIRECTORY")
openai_api_key = os.getenv("OPENAI_API_KEY")
openai.api_key = openai_api_key

# Extract text from files in the specified directory
def read_files_from_directory(directory):
    texts = []
    for filename in os.listdir(directory):
        file_path = os.path.join(directory, filename)
        text = ""

        try:
            if filename.endswith('.pdf'):
                with open(file_path, 'rb') as file:
                    reader = PdfReader(file)
                    text = " ".join(page.extract_text() or "" for page in reader.pages)
                texts.append(f"PDF: {filename}\n{text}")

            elif filename.endswith('.docx'):
                doc = Document(file_path)
                text = "\n".join([para.text for para in doc.paragraphs])
                texts.append(f"DOCX: {filename}\n{text}")

            elif filename.endswith('.pptx'):
                prs = Presentation(file_path)
                for slide in prs.slides:
                    slide_text = "\n".join([shape.text for shape in slide.shapes if hasattr(shape, "text")])
                    text += slide_text + "\n"
                texts.append(f"PPTX: {filename}\n{text}")

            elif filename.endswith(('.txt', '.md', '.py', '.c', '.cpp', '.java')):
                with open(file_path, 'r', encoding="utf-8", errors="ignore") as file:
                    text = file.read()
                texts.append(f"Text/Code: {filename}\n{text}")

            elif filename.endswith('.json'):
                with open(file_path, 'r', encoding="utf-8") as file:
                    data = json.load(file)
                    text = json.dumps(data, indent=2)
                texts.append(f"JSON: {filename}\n{text}")

            print(f"Processed: {filename}")

        except Exception as e:
            print(f"Error processing {filename}: {e}")
    
    return texts

# Split texts into chunks
def get_text_chunks(texts):
    text_splitter = RecursiveCharacterTextSplitter(chunk_size=3072, chunk_overlap=600)
    chunks = text_splitter.split_text(" ".join(texts))
    print("Chunks have been created.")
    return chunks

# Generate and save vector store with embeddings
def get_vector_store(text_chunks):
    embeddings = OpenAIEmbeddings(model="text-embedding-3-small")
    vector_store = FAISS.from_texts(text_chunks, embedding=embeddings)
    vector_store.save_local("faiss_index")
    print("Vector store has been saved locally as 'faiss_index'.")

# Save extracted texts or chunks to a file
def save_output_to_file(texts, filename):
    os.makedirs("output", exist_ok=True)
    file_path = os.path.join("output", filename)
    with open(file_path, "w", encoding="utf-8") as file:
        for text in texts:
            file.write(text + "\n\n")
    print(f"Output saved to {file_path}")

# Main function
def main():
    raw_texts = read_files_from_directory(directory)
    save_output_to_file(raw_texts, "extracted_texts.txt")
    
    text_chunks = get_text_chunks(raw_texts)
    save_output_to_file(text_chunks, "text_chunks.txt")
    
    get_vector_store(text_chunks)

if __name__ == "__main__":
    main()
